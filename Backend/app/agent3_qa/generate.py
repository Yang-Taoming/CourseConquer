"""多模态生成 skill：用统一模型 gpt-5.1-high（文本）+ gpt-image-1（图像）生成可下载文件。

每个生成类型是一个 skill 函数（不切换模型），返回 (文件绝对路径, 下载文件名, mime)。
生成内容尽量结合知识库：若给出 workspace_id，先检索相关证据作为生成素材。
"""
from __future__ import annotations

import io
import tempfile
from pathlib import Path
from typing import Optional, Tuple

from app.config import get_settings
from app.shared.llm import client as llm
from app.shared.storage.local import get_storage

# 生成类型 → (后缀, mime)
TYPE_EXT = {
    "notes": (".md", "text/markdown"),
    "report": (".md", "text/markdown"),
    "md": (".md", "text/markdown"),
    "ppt": (".pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    "doc": (".docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    "code": (".py", "text/x-python"),
    "csv": (".csv", "text/csv"),
    "image": (".png", "image/png"),
}


def _kb_context(workspace_id: Optional[str], topic: str) -> str:
    """从知识库检索与主题相关的证据，作为生成素材。"""
    if not workspace_id:
        return ""
    try:
        store = get_storage()
        docs = store.list_documents(workspace_id)
        if not docs:
            return ""
        emb = llm.embed([topic])[0]
        hits = store.search(workspace_id, emb, 4)
        if not hits:
            return ""
        ctx = "\n\n".join(h.get("text", "")[:500] for h in hits[:4])
        return "\n\n【知识库相关素材】\n" + ctx
    except Exception:
        return ""


def _gen_text(prompt: str, max_tokens: int = 4000) -> str:
    s = get_settings()
    return llm.chat(
        [{"role": "system", "content": "你是课程知识库的生成助手，按要求产出高质量中文内容。"},
         {"role": "user", "content": prompt}],
        model=s.llm_model, max_tokens=max_tokens,
    )


def classify_intent(prompt: str) -> Dict:
    """从用户自然语言判断要生成的产物类型（意图驱动）。"""
    s = get_settings()
    sys = (
        "你是生成意图分类器。根据用户的自然语言请求，判断要生成哪种产物。只输出 JSON："
        '{"kind":"image|ppt|doc|code|notes|report","topic":"生成主题（简短，去掉『画一张/做一个』等指令词）",'
        '"lang":"代码语言，仅 kind=code 时填，默认 python"}。'
        "规则：画图/示意图/插图/流程图→image；PPT/幻灯片/演示文稿→ppt；"
        "正式文档/报告/word→report；普通文档/说明→doc；代码/函数/实现/脚本→code；笔记/总结/复习提纲→notes。"
        "无法判断时默认 notes。"
    )
    d = llm.chat_json(sys, prompt[:500], model=s.llm_model)
    kind = d.get("kind", "notes")
    if kind not in ("image", "ppt", "doc", "code", "notes", "report"):
        kind = "notes"
    return {"kind": kind, "topic": str(d.get("topic", prompt[:24])).strip() or prompt[:24],
            "lang": str(d.get("lang", "python")).strip() or "python"}


def _write_temp(content: bytes, suffix: str) -> Path:
    fd, p = tempfile.mkstemp(suffix=suffix, prefix="gen_", dir="/tmp")
    with open(fd, "wb") as f:
        f.write(content)
    return Path(p)


def generate(kind: str, topic: str, workspace_id: Optional[str] = None,
             lang: str = "python") -> Tuple[Path, str, str]:
    """主入口：返回 (文件路径, 下载文件名, mime)。kind=auto 时按用户意图自动判别。"""
    s = get_settings()
    # 意图驱动：auto 时用 LLM 判别类型 + 提取主题
    if kind == "auto":
        intent = classify_intent(topic)
        kind = intent["kind"]
        topic = intent["topic"]
        lang = intent["lang"]
    ctx = _kb_context(workspace_id, topic)
    safe = "".join(c for c in topic[:24] if c.isalnum() or c in "一二三四五六七八九十") or "output"

    if kind == "image":
        img = llm.generate_image(topic + ctx[:200])
        suffix, mime = TYPE_EXT["image"]
        p = _write_temp(img, suffix)
        return p, f"{safe}.png", mime

    if kind == "code":
        prompt = (
            "请生成一份完整可运行的 %s 代码实现：'%s'。要求：只输出代码本身（带必要注释），"
            "不要解释文字、不要 markdown 围栏。%s" % (lang, topic, ctx)
        )
        code = _gen_text(prompt, 3000)
        ext = "." + lang if not lang.startswith(".") else lang
        mime = "text/plain"
        p = _write_temp(code.encode("utf-8"), ext)
        return p, f"{safe}{ext}", mime

    if kind == "csv":
        prompt = (
            "为主题 '%s' 生成一份示例 CSV 数据：第一行为表头，后续若干行数据。"
            "只输出 CSV 文本（逗号分隔），不要 markdown 围栏、不要解释。%s" % (topic, ctx)
        )
        content = _gen_text(prompt, 4000)
        p = _write_temp(content.encode("utf-8"), ".csv")
        return p, f"{safe}.csv", "text/csv"

    if kind == "md":
        prompt = "为主题 '%s' 生成一份 Markdown 文档（含标题、章节、要点）。%s" % (topic, ctx)
        content = _gen_text(prompt, 4000)
        p = _write_temp(content.encode("utf-8"), ".md")
        return p, f"{safe}.md", "text/markdown"

    if kind == "notes":
        prompt = "为以下主题生成结构化学习笔记（Markdown）：要点、定义、原理、示例、易错点。主题：%s%s" % (topic, ctx)
        content = _gen_text(prompt, 4000)
        p = _write_temp(content.encode("utf-8"), ".md")
        return p, f"{safe}_笔记.md", "text/markdown"

    if kind == "report":
        prompt = "为以下主题生成一份技术报告草稿（Markdown）：背景、方法、分析、结论。主题：%s%s" % (topic, ctx)
        content = _gen_text(prompt, 4000)
        p = _write_temp(content.encode("utf-8"), ".md")
        return p, f"{safe}_报告.md", "text/markdown"

    if kind == "ppt":
        prompt = (
            "为主题 '%s' 生成 PPT 大纲，输出严格 JSON：{\"slides\":[{\"title\":\"...\",\"bullets\":[\"...\"]}]}, 5-8 页。%s"
            % (topic, ctx)
        )
        import json
        d = llm.chat_json("你是 PPT 大纲生成器，只输出 JSON。", prompt, model=s.llm_model)
        slides = d.get("slides", [])
        from pptx import Presentation
        prs = Presentation()
        for sl in slides:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = str(sl.get("title", ""))
            body = slide.placeholders[1]
            for b in (sl.get("bullets") or [])[:6]:
                body.text = (body.text + "\n" if body.text else "") + str(b)
        buf = io.BytesIO()
        prs.save(buf)
        p = _write_temp(buf.getvalue(), ".pptx")
        return p, f"{safe}.pptx", TYPE_EXT["ppt"][1]

    if kind == "doc":
        prompt = "为主题 '%s' 生成一份文档（Markdown 段落，含标题与正文）。%s" % (topic, ctx)
        content = _gen_text(prompt, 4000)
        from docx import Document
        doc = Document()
        for line in content.splitlines():
            line = line.strip()
            if line.startswith("# "):
                doc.add_heading(line[2:], level=1)
            elif line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line:
                doc.add_paragraph(line)
        buf = io.BytesIO()
        doc.save(buf)
        p = _write_temp(buf.getvalue(), ".docx")
        return p, f"{safe}.docx", TYPE_EXT["doc"][1]

    raise ValueError("不支持的生成类型: %s" % kind)
