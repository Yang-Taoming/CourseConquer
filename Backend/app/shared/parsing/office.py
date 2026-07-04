"""Office / 表格：docx、pptx、xlsx、csv -> 归一化 Markdown。"""
from __future__ import annotations

import csv as csvmod
import io
from typing import List

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from app.shared.schemas.document import Block, BlockType, ParsedDocument

MAX_TABLE_ROWS = 200


def _table_to_md(rows: List[List[str]]) -> str:
    rows = [r for r in rows if any((c or "").strip() for c in r)]
    if not rows:
        return ""
    ncol = max(len(r) for r in rows)

    def norm(r: List[str]) -> List[str]:
        cells = []
        for i in range(ncol):
            v = r[i] if i < len(r) else ""
            cells.append((v or "").replace("\n", " ").replace("|", "\\|").strip())
        return cells

    header = norm(rows[0])
    out = ["| " + " | ".join(header) + " |",
           "| " + " | ".join(["---"] * ncol) + " |"]
    for r in rows[1:]:
        out.append("| " + " | ".join(norm(r)) + " |")
    return "\n".join(out)


def parse_docx(data: bytes, filename: str, mime: str) -> ParsedDocument:
    doc = DocxDocument(io.BytesIO(data))
    lines: List[str] = []
    blocks: List[Block] = []

    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading") or style == "title":
            lines.append("## " + t)
            blocks.append(Block(type=BlockType.heading, text=t))
        else:
            lines.append(t)
            blocks.append(Block(type=BlockType.paragraph, text=t))

    for table in doc.tables:
        md = _table_to_md([[c.text for c in row.cells] for row in table.rows])
        if md:
            lines.append(md)
            blocks.append(Block(type=BlockType.table, text=md))

    return ParsedDocument(
        filename=filename, mime=mime, doc_type="docx",
        markdown="\n\n".join(lines), blocks=blocks,
        meta={"parse_method": "python-docx"},
    )


def parse_pptx(data: bytes, filename: str, mime: str) -> ParsedDocument:
    prs = Presentation(io.BytesIO(data))
    lines: List[str] = []
    blocks: List[Block] = []
    n_slides = 0

    for idx, slide in enumerate(prs.slides, start=1):
        n_slides = idx
        lines.append("## Slide %d" % idx)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip() or para.text.strip()
                    if t:
                        lines.append(t)
                        blocks.append(Block(type=BlockType.paragraph, text=t, meta={"slide": idx}))
            if shape.has_table:
                md = _table_to_md([[c.text for c in row.cells] for row in shape.table.rows])
                if md:
                    lines.append(md)
                    blocks.append(Block(type=BlockType.table, text=md, meta={"slide": idx}))

    return ParsedDocument(
        filename=filename, mime=mime, doc_type="pptx",
        markdown="\n\n".join(lines), blocks=blocks,
        meta={"parse_method": "python-pptx", "slides": n_slides},
    )


def parse_xlsx(data: bytes, filename: str, mime: str) -> ParsedDocument:
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: List[str] = []
    blocks: List[Block] = []

    for ws in wb.worksheets:
        rows: List[List[str]] = []
        for r in ws.iter_rows(values_only=True):
            rows.append(["" if v is None else str(v) for v in r])
            if len(rows) >= MAX_TABLE_ROWS:
                break
        md = _table_to_md(rows)
        if md:
            lines.append("## " + str(ws.title))
            lines.append(md)
            blocks.append(Block(type=BlockType.table, text=md, meta={"sheet": ws.title}))
    wb.close()

    return ParsedDocument(
        filename=filename, mime=mime, doc_type="xlsx",
        markdown="\n\n".join(lines), blocks=blocks,
        meta={"parse_method": "openpyxl"},
    )


def parse_csv(data: bytes, filename: str, mime: str) -> ParsedDocument:
    text = data.decode("utf-8", errors="replace")
    rows = list(csvmod.reader(io.StringIO(text)))
    md = _table_to_md(rows[:MAX_TABLE_ROWS])
    blocks = [Block(type=BlockType.table, text=md)] if md else []
    return ParsedDocument(
        filename=filename, mime=mime, doc_type="csv",
        markdown=md, blocks=blocks,
        meta={"parse_method": "csv"},
    )
